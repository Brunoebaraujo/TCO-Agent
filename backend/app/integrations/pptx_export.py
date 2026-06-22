"""
pptx_export.py — Goodpack TCO Agent
Gera um .pptx branded a partir de um resultado de TCO estruturado,
usando o template oficial da Goodpack como base.

Compatível com o retorno de calculate_tco() em engine.py.
O template deve estar no mesmo diretório que este arquivo:
  app/integrations/230214_GP_Blank_Presentation_Template.pptx

Assinatura pública:
    generate_tco_pptx(tco: dict) -> bytes
    (retorna bytes prontos para resposta HTTP — sem gravar em disco)
"""

import io
import os
import math
import tempfile
from datetime import datetime
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paleta oficial Goodpack ──────────────────────────────────────────────────
GP_BLUE      = RGBColor(0x00, 0x4C, 0x97)
GP_BLUE_MID  = RGBColor(0x00, 0x5F, 0x9E)
GP_GREEN     = RGBColor(0x23, 0xAE, 0x49)
GP_ORANGE    = RGBColor(0xF2, 0x6B, 0x43)
GP_GREY      = RGBColor(0xA0, 0xA0, 0xA0)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)
CARD_BG      = RGBColor(0xF4, 0xF7, 0xFB)

HEX_BLUE     = "#004C97"
HEX_ORANGE   = "#F26B43"
HEX_GREEN    = "#23AE49"

TEMPLATE_PATH = Path(__file__).parent / "230214_GP_Blank_Presentation_Template.pptx"


# ── Helpers ──────────────────────────────────────────────────────────────────

def _num(v, default=0.0):
    return float(v) if v is not None else default


def _fmt(value, currency="USD", decimals=2):
    sym = "$" if currency == "USD" else f"{currency} "
    return f"{sym}{value:,.{decimals}f}"


def _textbox(slide, left, top, width, height, text,
             size=12, bold=False, italic=False,
             color=BLACK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _card(slide, left, top, width, height, border_color):
    rect = slide.shapes.add_shape(1, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = CARD_BG
    rect.line.color.rgb = border_color
    rect.line.width = Pt(1.5)
    bar = slide.shapes.add_shape(1, left, top, width, Inches(0.07))
    bar.fill.solid()
    bar.fill.fore_color.rgb = border_color
    bar.line.fill.background()


# ── Gráfico matplotlib ───────────────────────────────────────────────────────

def _make_chart(tco: dict) -> io.BytesIO:
    """2 barras: Goodpack vs Concorrente — apenas Packaging + Freight."""
    sub   = tco.get("subtotals", {})
    gp_v  = _num((sub.get("goodpack") or {}).get("packaging_and_freight"))
    cmp_v = _num((sub.get("competitor") or {}).get("packaging_and_freight"))

    # Fallback: se subtotals não existirem ainda, usa totais
    if not gp_v:
        gp_v  = _num(tco.get("goodpack_total_per_mt"))
    if not cmp_v:
        cmp_v = _num(tco.get("competitor_total_per_mt"))

    cur = tco.get("currency", "USD")
    sym = "$" if cur == "USD" else f"{cur} "
    gp_label  = tco.get("goodpack_sku", "Goodpack")
    cmp_label = tco.get("competitor_name", "Competitor")

    labels = [f"Goodpack\n{gp_label}", f"Competitor\n{cmp_label}"]
    values = [gp_v, cmp_v]

    fig, ax = plt.subplots(figsize=(6.2, 3.8), dpi=150)
    fig.patch.set_facecolor("white")
    ax.set_facecolor("white")

    bars = ax.bar(labels, values, color=[HEX_BLUE, HEX_ORANGE],
                  width=0.45, zorder=3, edgecolor="white", linewidth=0.5)

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2,
                bar.get_height() + max(values) * 0.015,
                f"{sym}{val:.2f}/MT",
                ha="center", va="bottom", fontsize=11,
                fontweight="bold", color="#1A1A1A")

    if cmp_v > 0:
        saving_pct = (cmp_v - gp_v) / cmp_v * 100
        ax.annotate(f"−{saving_pct:.1f}%\nsaving",
                    xy=(0.5, max(values) * 0.50),
                    xycoords=("axes fraction", "data"),
                    fontsize=9, ha="center",
                    color=HEX_GREEN, fontweight="bold")

    ax.set_ylabel(f"Cost per MT ({cur})", fontsize=10, color="#555555")
    ax.set_ylim(0, max(values) * 1.22)
    ax.yaxis.set_major_formatter(
        plt.FuncFormatter(lambda x, _: f"{sym}{x:.0f}"))
    ax.tick_params(axis="x", labelsize=11)
    ax.tick_params(axis="y", labelsize=9, colors="#888888")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.spines["bottom"].set_color("#DDDDDD")
    ax.yaxis.grid(True, linestyle="--", alpha=0.5, zorder=0)
    ax.set_axisbelow(True)

    fig.text(0.5, 0.01,
             "* Packaging & Freight only. Handling savings calculated separately.",
             ha="center", fontsize=7, color="#888888", style="italic")

    plt.tight_layout(rect=[0, 0.04, 1, 1])
    buf = io.BytesIO()
    plt.savefig(buf, format="png", dpi=150, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    buf.seek(0)
    return buf


# ── Slides ───────────────────────────────────────────────────────────────────

def _slide_cover(prs, tco):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    date_str = tco.get("date", datetime.today().strftime("%B %Y"))
    gp_pkg  = tco.get("goodpack_sku", "Goodpack")
    cmp_pkg = tco.get("competitor_name", "Competitor")

    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            ph.text = tco.get("customer_name", "Client")
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.bold = True
                    run.font.size = Pt(36)
                    run.font.color.rgb = GP_BLUE
        elif idx == 1:
            ph.text = f"TCO Analysis — {cmp_pkg} vs {gp_pkg}  |  {date_str}"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"
                    run.font.size = Pt(16)
                    run.font.color.rgb = GP_GREEN


def _slide_executive_summary(prs, tco):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    W, H  = prs.slide_width, prs.slide_height
    cur   = tco.get("currency", "USD")

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Executive Summary"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"; run.font.bold = True
                    run.font.size = Pt(28); run.font.color.rgb = GP_BLUE
        elif ph.placeholder_format.idx == 1:
            ph.text = ""

    total_saving    = _num(tco.get("total_saving"))
    saving_pct      = _num(tco.get("saving_percentage"))
    gp_total        = _num(tco.get("goodpack_total_per_mt"))
    mt              = _num(tco.get("simulated_metric_tonnes"))
    product         = tco.get("product_name", "")
    competitor_name = tco.get("competitor_name", "Competitor")

    kpis = [
        {"label": "Annual Savings",   "value": _fmt(total_saving, cur, 0),  "sub": f"vs {competitor_name}", "color": GP_GREEN},
        {"label": "Saving per MT",    "value": _fmt(saving_pct, "USD", 1).replace("$", "") + "%",  "sub": f"{_fmt(gp_total, cur)}/MT Goodpack", "color": GP_BLUE},
        {"label": "Annual Volume",    "value": f"{mt:,.0f} MT",              "sub": product,        "color": GP_ORANGE},
    ]
    # Fix saving per MT display
    saving_per_mt = _num(tco.get("competitor_total_per_mt")) - gp_total
    kpis[1]["value"] = _fmt(saving_per_mt, cur)
    kpis[1]["sub"]   = f"{saving_pct:.1f}% reduction"

    card_w = Inches(3.2); card_h = Inches(2.4); gap = Inches(0.35)
    total_w = 3 * card_w + 2 * gap
    start_x = (W - total_w) // 2
    card_y  = Inches(2.1)

    for i, kpi in enumerate(kpis):
        cx = start_x + i * (card_w + gap)
        _card(slide, cx, card_y, card_w, card_h, kpi["color"])
        _textbox(slide, cx + Inches(0.15), card_y + Inches(0.18),
                 card_w - Inches(0.3), Inches(0.9),
                 kpi["value"], size=28, bold=True,
                 color=kpi["color"], align=PP_ALIGN.CENTER)
        _textbox(slide, cx + Inches(0.1), card_y + Inches(1.05),
                 card_w - Inches(0.2), Inches(0.45),
                 kpi["label"], size=13, color=GP_BLUE, align=PP_ALIGN.CENTER)
        _textbox(slide, cx + Inches(0.1), card_y + Inches(1.55),
                 card_w - Inches(0.2), Inches(0.5),
                 kpi["sub"], size=11, color=GP_GREY,
                 align=PP_ALIGN.CENTER, italic=True)

    _textbox(slide, Inches(0.6), H - Inches(0.85), W - Inches(1.2), Inches(0.4),
             f"Based on {mt:,.0f} MT/year  |  Product: {product}  |  Packaging & Freight costs",
             size=9, color=GP_GREY, align=PP_ALIGN.CENTER, italic=True)


def _slide_chart(prs, tco):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    W, H  = prs.slide_width, prs.slide_height
    cur   = tco.get("currency", "USD")

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "TCO Breakdown: Packaging & Freight vs Full Cost"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"; run.font.bold = True
                    run.font.size = Pt(24); run.font.color.rgb = GP_BLUE
        elif ph.placeholder_format.idx == 1:
            ph.text = ""

    chart_buf = _make_chart(tco)
    slide.shapes.add_picture(chart_buf, Inches(1.2), Inches(1.5),
                             Inches(6.8), Inches(4.3))

    panel_x = Inches(8.4); panel_y = Inches(1.6); panel_w = Inches(4.2)
    gp_total  = _num(tco.get("goodpack_total_per_mt"))
    cmp_total = _num(tco.get("competitor_total_per_mt"))
    saving_mt = cmp_total - gp_total
    saving_an = _num(tco.get("total_saving"))

    kpis = [
        ("Goodpack total/MT",    _fmt(gp_total,  cur), GP_BLUE),
        ("Competitor total/MT",  _fmt(cmp_total, cur), GP_ORANGE),
        ("Saving/MT (full TCO)", _fmt(saving_mt, cur), GP_GREEN),
        ("Annual saving",        _fmt(saving_an, cur, 0), GP_GREEN),
    ]
    for i, (label, val, color) in enumerate(kpis):
        row_y = panel_y + i * Inches(0.95)
        _card(slide, panel_x, row_y, panel_w, Inches(0.80), color)
        _textbox(slide, panel_x + Inches(0.12), row_y + Inches(0.04),
                 panel_w - Inches(0.15), Inches(0.32), label, size=9, color=GP_GREY)
        _textbox(slide, panel_x + Inches(0.12), row_y + Inches(0.35),
                 panel_w - Inches(0.15), Inches(0.38), val,
                 size=14, bold=True, color=color)


def _slide_detail_cards(prs, tco):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    W, H  = prs.slide_width, prs.slide_height
    cur   = tco.get("currency", "USD")

    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = "Cost Breakdown by Category"
            for para in ph.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = "Calibri"; run.font.bold = True
                    run.font.size = Pt(28); run.font.color.rgb = GP_BLUE
        elif ph.placeholder_format.idx == 1:
            ph.text = ""

    categories = tco.get("categories", [])
    if not categories:
        _textbox(slide, Inches(1), Inches(2), Inches(10), Inches(1),
                 "No category data available.", size=14, color=GP_GREY)
        return

    n      = len(categories)
    cols   = min(n, 3)
    card_w = Inches(3.6); card_h = Inches(3.0); gap = Inches(0.35)
    total_w = cols * card_w + (cols - 1) * gap
    start_x = (W - total_w) // 2
    card_y  = Inches(1.9)

    for i, cat in enumerate(categories):
        col = i % 3; row = i // 3
        cx  = start_x + col * (card_w + gap)
        cy  = card_y + row * (card_h + Inches(0.2))

        gp_v    = _num(cat.get("goodpack"))
        cmp_v   = _num(cat.get("competitor"))
        diff    = cmp_v - gp_v
        pct     = (diff / cmp_v * 100) if cmp_v else 0
        color   = GP_GREEN if pct >= 0 else GP_ORANGE

        _card(slide, cx, cy, card_w, card_h, color)
        _textbox(slide, cx + Inches(0.15), cy + Inches(0.12),
                 card_w - Inches(0.3), Inches(0.4),
                 cat.get("label", cat.get("name", "")).upper(),
                 size=12, bold=True, color=GP_BLUE, align=PP_ALIGN.CENTER)
        sign = "−" if pct >= 0 else "+"
        _textbox(slide, cx + Inches(0.1), cy + Inches(0.55),
                 card_w - Inches(0.2), Inches(0.75),
                 f"{sign}{abs(pct):.1f}%",
                 size=32, bold=True, color=color, align=PP_ALIGN.CENTER)
        _textbox(slide, cx + Inches(0.1), cy + Inches(1.28),
                 card_w - Inches(0.2), Inches(0.3),
                 "saving vs competitor",
                 size=9, color=GP_GREY, align=PP_ALIGN.CENTER, italic=True)
        _textbox(slide, cx + Inches(0.15), cy + Inches(1.7),
                 card_w - Inches(0.3), Inches(0.32),
                 f"Goodpack:    {_fmt(gp_v, cur)}/MT",
                 size=10, bold=True, color=GP_BLUE)
        _textbox(slide, cx + Inches(0.15), cy + Inches(2.05),
                 card_w - Inches(0.3), Inches(0.32),
                 f"Competitor:  {_fmt(cmp_v, cur)}/MT",
                 size=10, color=GP_ORANGE)
        _textbox(slide, cx + Inches(0.15), cy + Inches(2.4),
                 card_w - Inches(0.3), Inches(0.32),
                 f"Saving:      {_fmt(diff, cur)}/MT",
                 size=10, bold=True, color=color)


def _slide_closing(prs, tco):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    W, H  = prs.slide_width, prs.slide_height

    for ph in slide.placeholders:
        try:
            ph.text = ""
        except Exception:
            pass

    # Tagline bicolor
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(2.6),
                                  W - Inches(1.0), Inches(1.0))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    for text, color, bold in [
        ("CREATING ",       GP_BLUE,  False),
        ("CIRCULAR",        GP_GREEN, True),
        (" SUPPLY CHAINS",  GP_BLUE,  False),
    ]:
        r = p.add_run()
        r.text = text
        r.font.name = "Calibri"; r.font.size = Pt(30)
        r.font.bold = bold; r.font.italic = True
        r.font.color.rgb = color

    _textbox(slide, Inches(0.5), Inches(3.85), W - Inches(1.0), Inches(0.5),
             "www.goodpack.com", size=14, color=GP_BLUE_MID,
             align=PP_ALIGN.CENTER)
    _textbox(slide, Inches(0.5), H - Inches(1.35), W - Inches(1.0), Inches(0.45),
             "Copyright © 2023 Goodpack IBC (Singapore) Pte Ltd. All rights reserved. "
             "GOODPACK, TYRECUBE and TOMATOCUBE are registered trademarks of "
             "Goodpack IBC (Singapore) Pte Ltd. Technical and commercial information "
             "in this document is subject to change without notice.",
             size=7, color=GP_GREY, align=PP_ALIGN.CENTER)


# ── API pública ──────────────────────────────────────────────────────────────

def generate_tco_pptx(tco: dict, template_path: str = None,
                       include_assumptions: bool = True) -> bytes:
    """
    Gera o .pptx brandado da Goodpack e retorna bytes prontos para HTTP.

    Parâmetro `include_assumptions` mantido por compatibilidade com o
    endpoint existente em tco.py — não é usado nesta versão (o slide de
    premissas foi substituído pelo slide de cards por categoria).
    """
    tmpl = Path(template_path) if template_path else TEMPLATE_PATH
    if not tmpl.exists():
        raise FileNotFoundError(
            f"Template Goodpack não encontrado: {tmpl}\n"
            "Coloque 230214_GP_Blank_Presentation_Template.pptx em "
            "backend/app/integrations/"
        )

    prs = Presentation(str(tmpl))

    # Remove os slides placeholder do template em branco
    for i in range(len(prs.slides) - 1, -1, -1):
        rId = prs.slides._sldIdLst[i].get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[i])

    _slide_cover(prs, tco)
    _slide_executive_summary(prs, tco)
    _slide_chart(prs, tco)
    _slide_detail_cards(prs, tco)
    _slide_closing(prs, tco)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
